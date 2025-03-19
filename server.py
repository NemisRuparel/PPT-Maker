from flask import Flask, request, send_file
from pptx import Presentation
import os
import docx  # For reading .docx files (install with `pip install python-docx`)

app = Flask(__name__)

@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    reference_file = request.files.get("reference")
    template_file = request.files.get("template")

    if not reference_file:
        return {"error": "Reference file is required"}, 400

    # Save files temporarily
    ref_path = "temp_reference" + os.path.splitext(reference_file.filename)[1]
    reference_file.save(ref_path)
    template_path = template_file.save("temp_template.pptx") if template_file else None

    # Extract text from reference file (assuming .docx for simplicity)
    content = ""
    if ref_path.endswith(".docx"):
        doc = docx.Document(ref_path)
        content = "\n".join([para.text for para in doc.paragraphs])
    else:
        with open(ref_path, "r", encoding="utf-8") as f:
            content = f.read()

    # Create PPT using template or new presentation
    prs = Presentation(template_path) if template_path else Presentation()
    
    # Simple slide generation: split content into slides
    lines = content.split("\n")
    for i in range(0, len(lines), 3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = f"Slide {i // 3 + 1}"
        body.text = "\n".join(lines[i:i+3])

    # Save the PPTX
    output_path = "generated_slides.pptx"
    prs.save(output_path)

    # Clean up temporary files
    os.remove(ref_path)
    if template_path:
        os.remove(template_path)

    # Return the PPTX file
    return send_file(output_path, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)